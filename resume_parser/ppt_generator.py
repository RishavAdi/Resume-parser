from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.table import _Cell
from typing import Dict, Any

class ResumeSlideGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

        # Define consistent styling
        self.styles = {
            'header': {
                'font_size': Pt(16),
                'font_color': RGBColor(47, 84, 150),
                'bold': True
            },
            'content': {
                'font_size': Pt(11),
                'font_color': RGBColor(0, 0, 0),
                'bold': False
            }
        }

    def _apply_cell_style(self, cell: _Cell, is_header: bool = False):
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        font = paragraph.font
        style = self.styles['header'] if is_header else self.styles['content']
        font.size = style['font_size']
        font.color.rgb = style['font_color']
        font.bold = style['bold']

        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.1)
        cell.margin_top = cell.margin_bottom = Inches(0.05)

    def generate_slide(self, resume_data: Dict[str, Any], output_path: str) -> str:
        try:
            slide_layout = self.prs.slide_layouts[6]
            slide = self.prs.slides.add_slide(slide_layout)

            # Define layout constants
            left_margin = Inches(0.5)
            right_margin = Inches(0.5)
            column_gap = Inches(0.5)

            # Total usable width (EMU ints)
            content_width = int(self.slide_width - (left_margin + right_margin + column_gap))
            # Use integer division to avoid floats
            column_width = int(content_width // 2)

            # Convert all positions to int (EMU)
            left_col_x = int(left_margin)
            right_col_x = int(left_margin + column_width + column_gap)

            # 1️⃣ Personal Info (Top-left)
            personal_box = slide.shapes.add_textbox(
                left_col_x, int(Inches(0.4)), int(column_width), int(Inches(1))
            )
            self._add_personal_info(personal_box, resume_data)

            # 2️⃣ Headline Summary (Bottom-left)
            summary_box = slide.shapes.add_textbox(
                left_col_x, int(Inches(1.6)), int(column_width), int(Inches(3))
            )
            self._add_summary_section(summary_box, resume_data)

            # 3️⃣ Key Responsibilities (Top-right)
            resp_box = slide.shapes.add_textbox(
                right_col_x, int(Inches(0.4)), int(column_width), int(Inches(3))
            )
            self._add_responsibilities(resp_box, resume_data)

            # 4️⃣ Skills Table (Bottom-right)
            self._add_skills_table(slide, resume_data, right_col_x, int(Inches(3.7)), int(column_width))

            self.prs.save(output_path)
            return output_path

        except Exception as e:
            raise Exception(f"Failed to generate PowerPoint: {str(e)}")


    def _add_personal_info(self, textbox, data: Dict[str, Any]):
        tf = textbox.text_frame
        tf.clear()
        tf.word_wrap = True

        name_p = tf.paragraphs[0]
        name_p.text = data.get('name', 'N/A')
        name_p.font.bold = True
        name_p.font.size = Pt(18)
        name_p.font.color.rgb = RGBColor(47, 84, 150)

        contact_p = tf.add_paragraph()
        contact = data.get('email', '')
        if data.get('phone'):
            contact += f" | {data.get('phone')}"
        contact_p.text = contact
        contact_p.font.size = Pt(10)

        links = []
        if data.get('linkedin'):
            links.append(f"LinkedIn: {data['linkedin']}")
        if data.get('github'):
            links.append(f"GitHub: {data['github']}")
        if links:
            links_p = tf.add_paragraph()
            links_p.text = " | ".join(links)
            links_p.font.size = Pt(10)

    def _add_summary_section(self, textbox, data: Dict[str, Any]):
        tf = textbox.text_frame
        tf.clear()
        tf.word_wrap = True

        # Header
        header = tf.paragraphs[0]
        header.text = "Professional Summary"
        header.font.bold = True
        header.font.size = Pt(14)
        header.font.color.rgb = RGBColor(47, 84, 150)

        # Summary text
        if data.get('headline_summary'):
            p = tf.add_paragraph()
            p.text = data['headline_summary']
            p.font.size = Pt(11)

        # Certifications
        if data.get('certifications'):
            cert_header = tf.add_paragraph()
            cert_header.text = "\nCertifications:"
            cert_header.font.bold = True
            cert_header.font.size = Pt(12)
            for cert in data['certifications']:
                p = tf.add_paragraph()
                cert_text = f"• {cert.get('name', '')}"
                if cert.get('issuer'):
                    cert_text += f" ({cert.get('issuer')})"
                if cert.get('date'):
                    cert_text += f" - {cert.get('date')}"
                p.text = cert_text
                p.font.size = Pt(10)

    def _add_responsibilities(self, textbox, data: Dict[str, Any]):
        tf = textbox.text_frame
        tf.clear()
        tf.word_wrap = True

        header = tf.paragraphs[0]
        header.text = "Key Responsibilities"
        header.font.bold = True
        header.font.size = Pt(14)
        header.font.color.rgb = RGBColor(47, 84, 150)

        if data.get('key_responsibilities'):
            for resp in data['key_responsibilities']:
                p = tf.add_paragraph()
                p.text = f"• {resp}"
                p.font.size = Pt(10)
                p.space_after = Pt(2)

    def _add_skills_table(self, slide, data: Dict[str, Any], left: int, top: int, width: int):
        skills = data.get('skills', {})
        if not skills:
            return

        primary = skills.get('primary', [])
        secondary = skills.get('secondary', [])

        # Add header
        title_box = slide.shapes.add_textbox(int(left), int(top), int(width), int(Inches(0.3)))
        title_para = title_box.text_frame.paragraphs[0]
        title_para.text = "Technical Skills"
        title_para.font.bold = True
        title_para.font.size = Pt(14)
        title_para.font.color.rgb = RGBColor(47, 84, 150)

        rows = max(len(primary), len(secondary), 1) + 1
        # calculate table height in Inches then convert to EMU int
        table_height = int(Inches(0.3 * rows + 0.2))
        table = slide.shapes.add_table(rows, 2, int(left), int(top + Inches(0.4)), int(width), table_height).table

        # Ensure integer column widths
        col_w = int(width // 2)
        table.columns[0].width = col_w
        table.columns[1].width = col_w

        headers = ["Primary Skills", "Secondary Skills"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._apply_cell_style(cell, is_header=True)

        for i in range(1, rows):
            pi = i - 1
            for j, skill_list in enumerate([primary, secondary]):
                cell = table.cell(i, j)
                if pi < len(skill_list):
                    cell.text = str(skill_list[pi])
                else:
                    cell.text = ""
                self._apply_cell_style(cell)
